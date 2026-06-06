"""Inspect a .pptx template and extract layout information."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from ..schemas.template_config import TemplateConfig, LayoutMapping, Dimensions, SafeZone, BrandConfig


def profile_template(template_path: str) -> TemplateConfig:
    prs = Presentation(template_path)

    width_emu = prs.slide_width
    height_emu = prs.slide_height

    brand = _extract_brand(prs)
    layout_mappings = {}
    content_type_routing = {}

    for idx, layout in enumerate(prs.slide_layouts):
        layout_name = layout.name.lower().replace(" ", "_").replace("-", "_")
        placeholders = {}
        use_case = None

        for ph in layout.placeholders:
            ph_type = str(ph.placeholder_format.type)
            ph_name = ph.name.lower()

            if "title" in ph_name or "title" in ph_type.lower():
                placeholders["title"] = ph.placeholder_format.idx
                use_case = "title"
            elif "body" in ph_name or "object" in ph_type.lower():
                placeholders["body"] = ph.placeholder_format.idx
                use_case = "bullets"
            elif "subtitle" in ph_name:
                placeholders["subtitle"] = ph.placeholder_format.idx
            elif "picture" in ph_name:
                placeholders["image"] = ph.placeholder_format.idx

        layout_mappings[layout_name] = LayoutMapping(
            layout_index=idx,
            layout_name=layout_name,
            use_case=use_case,
        )

    for name, mapping in layout_mappings.items():
        if "title_scene" in name or "title" == name:
            content_type_routing["title"] = name
        elif "stat" in name or "big_number" in name:
            content_type_routing["stat"] = name
        elif "bullet" in name:
            content_type_routing["bullets"] = name
        elif "quote" in name:
            content_type_routing["quote"] = name
        elif "image" in name:
            content_type_routing["image_with_caption"] = name
        elif "cta" in name:
            content_type_routing["cta"] = name
        elif "list" in name or "numbered" in name:
            content_type_routing["numbered_list"] = name
        elif "section" in name:
            content_type_routing["section"] = name

    fallback_idx = 0
    fallback_name = list(layout_mappings.keys())[0] if layout_mappings else "bullet_scene"

    return TemplateConfig(
        template_name=Path(template_path).stem,
        template_path=template_path,
        dimensions=Dimensions(width_emu=width_emu, height_emu=height_emu),
        safe_zone=SafeZone(),
        brand=brand,
        layout_mappings=layout_mappings,
        content_type_routing=content_type_routing,
        fallback_layout_index=fallback_idx,
        fallback_layout_name=fallback_name,
    )


def _extract_brand(prs: Presentation) -> BrandConfig:
    primary_color = "0196FF"
    header_font = "Aptos"
    body_font = "Aptos"

    try:
        theme = prs.slide_masters[0].slide_layouts[0]
        for ph in theme.placeholders:
            if ph.has_text_frame:
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            header_font = run.font.name
                            body_font = run.font.name
                            break
                break
    except Exception:
        pass

    return BrandConfig(
        primary_color=primary_color,
        header_font=header_font,
        body_font=body_font,
    )