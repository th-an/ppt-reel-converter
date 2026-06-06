"""Fill template placeholders with content. NEVER add freeform text boxes.""" ""

from __future__ import annotations

from pptx.slide import Slide


def fill_placeholder(
    slide: Slide,
    placeholder_type: str,
    content: str,
) -> bool:
    for ph in slide.placeholders:
        ph_type = str(ph.placeholder_format.type)
        if placeholder_type == "title" and ph_type in ("TITLE", "CENTER_TITLE"):
            ph.text = content
            return True
        elif placeholder_type == "body" and ph_type in ("BODY", "OBJECT"):
            ph.text = content
            return True
        elif placeholder_type == "subtitle" and ph_type == "SUBTITLE":
            ph.text = content
            return True
    for ph in slide.placeholders:
        if ph.has_text_frame:
            ph.text_frame.text = content
            return True
    return False


def discover_placeholders(slide: Slide) -> dict[str, dict]:
    placeholders = {}
    for ph in slide.placeholders:
        ph_type = str(ph.placeholder_format.type)
        placeholders[ph.placeholder_format.idx] = {
            "index": ph.placeholder_format.idx,
            "type": ph_type,
            "name": ph.name,
            "left": ph.left,
            "top": ph.top,
            "width": ph.width,
            "height": ph.height,
        }
    return placeholders