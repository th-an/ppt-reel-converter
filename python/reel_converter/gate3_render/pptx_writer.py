"""Write actual PPTX slides using template layouts and real placeholder filling."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from ..schemas.scene_plan import ScenePlan, Scene
from ..schemas.template_config import TemplateConfig
from ..schemas.render_result import RenderedScene


def write_scenes_to_pptx(
    scene_plan: ScenePlan,
    template_config: TemplateConfig,
    output_path: str,
    original_images: dict[str, str] | None = None,
) -> list[RenderedScene]:
    """Generate actual PPTX file from scene plan using template layouts.
    
    Args:
        scene_plan: The planned scenes
        template_config: Template configuration with layout mappings
        output_path: Where to save the output PPTX
        original_images: Dict of image_name -> file_path for images to insert
    
    Returns:
        List of rendered scenes with metadata
    """
    template_path = template_config.template_path or "python/reel_converter/templates/reel_clean.pptx"
    prs = Presentation(template_path)
    
    # Remove all existing slides (they were just examples)
    slide_id_list = prs.slides._sldIdLst
    while len(slide_id_list) > 0:
        slide_id_entry = slide_id_list[0]
        slide_id = slide_id_entry.get(qn('r:id'))
        slide_id_list.remove(slide_id_entry)
        prs.part.drop_rel(slide_id)
    
    rendered = []
    
    for idx, scene in enumerate(scene_plan.scenes, start=1):
        rendered_scene = _write_scene(
            prs=prs,
            scene=scene,
            scene_number=idx,
            template_config=template_config,
            original_images=original_images,
        )
        rendered.append(rendered_scene)
    
    # Save output
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    
    return rendered


def _write_scene(
    prs: Presentation,
    scene: Scene,
    scene_number: int,
    template_config: TemplateConfig,
    original_images: dict[str, str] | None = None,
) -> RenderedScene:
    """Write a single scene to the presentation."""
    layout_name = scene.layout
    layout_mapping = template_config.layout_mappings.get(layout_name)
    
    if not layout_mapping:
        layout_mapping = template_config.layout_mappings.get(
            template_config.fallback_layout_name
        )
    
    layout_idx = layout_mapping.layout_index if layout_mapping else 0
    slide_layout = prs.slide_layouts[layout_idx]
    
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background color
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    background.line.fill.background()
    
    # Send background to back
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    
    placeholders_filled = {}
    images_inserted = []
    has_overflow = False
    min_font_size = None
    
    def _set_text(shape, text, font_size, font_color, bold=False, italic=False, align=PP_ALIGN.CENTER):
        """Clear text frame and add text with proper formatting."""
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = font_color
        return p
    
    # Fill placeholders based on scene type
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        
        ph = shape.placeholder_format
        ph_type = ph.type
        
        # TITLE placeholder (including CENTER_TITLE)
        if ph_type in (1, 3):  # TITLE or CENTER_TITLE
            if scene.stat_number:
                _set_text(shape, scene.stat_number, 72, RGBColor(0x01, 0x96, 0xFF), bold=True)
                placeholders_filled["title"] = scene.stat_number
            elif scene.headline:
                _set_text(shape, scene.headline, 48, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
                placeholders_filled["title"] = scene.headline
            elif scene.quote_text:
                _set_text(shape, f'"{scene.quote_text}"', 36, RGBColor(0xFF, 0xFF, 0xFF), italic=True)
                placeholders_filled["title"] = scene.quote_text
            elif scene.cta_headline:
                _set_text(shape, scene.cta_headline, 48, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
                placeholders_filled["title"] = scene.cta_headline
        
        # BODY / CONTENT placeholder
        elif ph_type in (2, 7):  # BODY or OBJECT
            if scene.body_items:
                tf = shape.text_frame
                tf.clear()
                for i, item in enumerate(scene.body_items):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(24)
                    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    p.space_after = Pt(12)
                
                text = "\n".join(f"• {item}" for item in scene.body_items)
                placeholders_filled["body"] = text
            elif scene.stat_label:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = scene.stat_label
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                if scene.stat_sublabel:
                    p2 = tf.add_paragraph()
                    p2.text = scene.stat_sublabel
                    p2.font.size = Pt(20)
                    p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
                
                label = scene.stat_label
                if scene.stat_sublabel:
                    label += f"\n{scene.stat_sublabel}"
                placeholders_filled["body"] = label
        
        # SUBTITLE placeholder
        elif ph_type == 4:  # SUBTITLE
            if scene.stat_sublabel:
                _set_text(shape, scene.stat_sublabel, 24, RGBColor(0xFF, 0xFF, 0xFF))
                placeholders_filled["subtitle"] = scene.stat_sublabel
            elif scene.quote_attribution:
                _set_text(shape, f"— {scene.quote_attribution}", 24, RGBColor(0xCC, 0xCC, 0xCC), italic=True)
                placeholders_filled["subtitle"] = scene.quote_attribution
            elif scene.cta_subheadline:
                _set_text(shape, scene.cta_subheadline, 24, RGBColor(0xFF, 0xFF, 0xFF))
                placeholders_filled["subtitle"] = scene.cta_subheadline
        
        # PICTURE placeholder
        elif ph_type == 18:  # PICTURE
            if scene.image_name and original_images:
                img_path = original_images.get(scene.image_name)
                if img_path and Path(img_path).exists():
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
                    images_inserted.append(scene.image_name)
                else:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
                    placeholders_filled["image"] = "placeholder"
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
                placeholders_filled["image"] = "placeholder"
    
    # If we have an image but no picture placeholder, add it as a shape
    if scene.image_name and original_images and not images_inserted:
        img_path = original_images.get(scene.image_name)
        if img_path and Path(img_path).exists():
            from PIL import Image as PILImage
            with PILImage.open(img_path) as img:
                img_w, img_h = img.size
            
            max_w = int(prs.slide_width * 0.9)
            max_h = int(prs.slide_height * 0.5)
            
            ratio = min(max_w / img_w, max_h / img_h)
            new_w = int(img_w * ratio)
            new_h = int(img_h * ratio)
            
            left = (prs.slide_width - new_w) // 2
            top = int(prs.slide_height * 0.25)
            
            slide.shapes.add_picture(img_path, left, top, new_w, new_h)
            images_inserted.append(scene.image_name)
    
    return RenderedScene(
        scene_number=scene_number,
        layout_used=layout_name,
        placeholders_filled=placeholders_filled,
        images_inserted=images_inserted,
        has_text_overflow=has_overflow,
        min_font_size_pt=min_font_size,
    )
