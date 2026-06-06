"""Visual validation using text overflow detection."""

from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu


def check_text_overflow(pptx_path: str) -> list[dict]:
    """Check for text overflow in all slides of a PPTX.
    
    Returns list of overflow detections with slide number, shape name, and severity.
    """
    prs = Presentation(pptx_path)
    overflows = []
    
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            tf = shape.text_frame
            text = tf.text.strip()
            if not text:
                continue
            
            # Get shape dimensions
            width = shape.width
            height = shape.height
            
            # Estimate text area
            # A rough heuristic: count characters and estimate based on font size
            total_chars = len(text)
            total_lines = len(text.split('\n'))
            
            # Get font size (use first run's font size or default 18pt)
            font_size = Pt(18)
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        font_size = run.font.size
                        break
                break
            
            # Estimate text dimensions
            # Rough estimate: char width ≈ font_size * 0.5, line height ≈ font_size * 1.2
            chars_per_line = int(width / (font_size * 0.5))
            estimated_lines = max(total_lines, int(total_chars / max(chars_per_line, 1)) + 1)
            estimated_height = estimated_lines * font_size * 1.2
            
            # Check if estimated height exceeds shape height
            if estimated_height > height:
                overflow_ratio = estimated_height / height
                severity = "high" if overflow_ratio > 1.5 else "medium" if overflow_ratio > 1.2 else "low"
                overflows.append({
                    "slide": slide_idx,
                    "shape": shape.name,
                    "text_preview": text[:50] + "..." if len(text) > 50 else text,
                    "estimated_height": int(estimated_height),
                    "shape_height": int(height),
                    "overflow_ratio": round(overflow_ratio, 2),
                    "severity": severity,
                })
    
    return overflows


def validate_pptx_visual(pptx_path: str) -> dict:
    """Perform visual validation on a PPTX file.
    
    Returns validation report with overflow detections and recommendations.
    """
    overflows = check_text_overflow(pptx_path)
    
    passed = len(overflows) == 0
    
    return {
        "passed": passed,
        "overflow_count": len(overflows),
        "overflows": overflows,
        "recommendations": [
            "Reduce text content" if o["severity"] == "high" else "Consider reducing text"
            for o in overflows
        ] if overflows else ["No overflow detected"],
    }
