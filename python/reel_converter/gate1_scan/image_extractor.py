"""Extract and save images from PPTX files."""

from __future__ import annotations

import os
import hashlib
from pathlib import Path
from typing import BinaryIO

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_images_from_pptx(pptx_path: str, output_dir: str) -> dict[str, str]:
    """Extract all images from a PPTX file.
    
    Args:
        pptx_path: Path to the PPTX file
        output_dir: Directory to save extracted images
    
    Returns:
        Dict mapping image names to file paths
    """
    prs = Presentation(pptx_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    
    images = {}
    image_count = 0
    
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_blob = image.blob
                    image_ext = image.ext
                    
                    # Create a unique filename
                    image_hash = hashlib.md5(image_blob).hexdigest()[:8]
                    image_count += 1
                    image_name = f"slide{slide_idx}_img{image_count}_{image_hash}.{image_ext}"
                    image_path = output_dir / image_name
                    
                    # Save image
                    with open(image_path, 'wb') as f:
                        f.write(image_blob)
                    
                    # Store with slide reference
                    key = f"slide_{slide_idx}_image_{image_count}"
                    images[key] = str(image_path)
                    
                except Exception as e:
                    print(f"Warning: Could not extract image from slide {slide_idx}: {e}")
    
    return images


def extract_slide_images(slide, slide_number: int, output_dir: str) -> dict[str, str]:
    """Extract images from a single slide.
    
    Args:
        slide: The slide object
        slide_number: The slide number
        output_dir: Directory to save extracted images
    
    Returns:
        Dict mapping image names to file paths
    """
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    
    images = {}
    image_count = 0
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = shape.image
                image_blob = image.blob
                image_ext = image.ext
                
                # Create a unique filename
                image_hash = hashlib.md5(image_blob).hexdigest()[:8]
                image_count += 1
                image_name = f"slide{slide_number}_img{image_count}_{image_hash}.{image_ext}"
                image_path = output_dir / image_name
                
                # Save image
                with open(image_path, 'wb') as f:
                    f.write(image_blob)
                
                key = f"image_{image_count}"
                images[key] = str(image_path)
                
            except Exception as e:
                print(f"Warning: Could not extract image from slide {slide_number}: {e}")
    
    return images
