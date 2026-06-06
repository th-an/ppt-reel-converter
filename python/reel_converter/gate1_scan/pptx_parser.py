"""Parse a single slide into structured elements."""

from __future__ import annotations

from pptx.slide import Slide


def parse_slide(slide: Slide, slide_height: int = 9144000) -> list[dict]:
    elements = []
    for shape in slide.shapes:
        el = _parse_shape(shape, slide_height)
        if el:
            elements.append(el)
    return elements


def _parse_shape(shape, slide_height: int = 9144000) -> dict | None:
    el = {
        "shape_id": shape.shape_id,
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
        "slide_height": slide_height,
        "rotation": getattr(shape, "rotation", 0),
    }

    if shape.has_text_frame:
        paragraphs = []
        all_text = []
        for para in shape.text_frame.paragraphs:
            runs_text = []
            for run in para.runs:
                    runs_text.append(run.text)
                    if run.font.name:
                        el["font_name"] = run.font.name
                    if run.font.size:
                        el["font_size"] = run.font.size
                    if run.font.bold:
                        el["font_bold"] = run.font.bold
                    if run.font.italic:
                        el["font_italic"] = run.font.italic
                    try:
                        if run.font.color and run.font.color.rgb:
                            el["font_color"] = str(run.font.color.rgb)
                    except AttributeError:
                        pass
            # Also check paragraph-level font if no run font
            if not el.get("font_size") and para.font.size:
                el["font_size"] = para.font.size
            if not el.get("font_name") and para.font.name:
                el["font_name"] = para.font.name
            para_text = "".join(runs_text).strip()
            if para_text:
                paragraphs.append({"text": para_text, "level": para.level})
                all_text.append(para_text)

        el["text"] = "\n".join(all_text)
        el["paragraphs"] = paragraphs

    if shape.shape_type is not None and "PICTURE" in str(shape.shape_type):
        el["is_image"] = True
        try:
            el["image_content_type"] = shape.image.content_type
            el["image_blob"] = shape.image.blob
        except Exception:
            pass

    if shape.has_table:
        table = shape.table
        rows = []
        for row_idx, row in enumerate(table.rows):
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip())
            rows.append(row_data)
        el["is_table"] = True
        el["table_rows"] = rows

    try:
        if shape.chart:
            el["is_chart"] = True
    except (ValueError, AttributeError):
        pass

    try:
        if shape.group_shapes:
            el["is_group"] = True
    except (ValueError, AttributeError):
        pass

    return el