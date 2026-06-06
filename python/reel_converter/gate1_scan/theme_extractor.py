"""Extract theme information from a PPTX presentation."""

from __future__ import annotations

from pptx.presentation import Presentation

from ..schemas.slide_model import Theme


def extract_theme(prs: Presentation) -> Theme:
    theme = Theme()
    try:
        slide_master = prs.slide_masters[0]
        if slide_master.slide_layouts:
            pass
    except Exception:
        pass
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name and not theme.header_font:
                            theme.header_font = run.font.name
                            theme.body_font = run.font.name
                        if run.font.color and run.font.color.rgb:
                            rgb = str(run.font.color.rgb)
                            if not theme.primary_color:
                                theme.primary_color = rgb
                            elif not theme.secondary_color:
                                theme.secondary_color = rgb
                            break
                    break
            break
        break
    return theme