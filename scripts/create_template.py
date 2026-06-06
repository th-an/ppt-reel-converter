from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Instagram Reels 9:16 aspect ratio
# Standard PowerPoint 9:16 = 5.625" x 10" (13.333" x 7.5" is landscape)
# For 9:16 portrait, we need width = 5.625, height = 10
# 1 inch = 914400 EMU
# 5.625" = 5143500 EMU
# 10" = 9144000 EMU

prs = Presentation()
prs.slide_width = Emu(5143500)
prs.slide_height = Emu(9144000)

# Color palette
PRIMARY_COLOR = RGBColor(0x01, 0x96, 0xFF)  # #0196FF
SECONDARY_COLOR = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)

# Safe zones
SAFE_TOP = 0.15
SAFE_BOTTOM = 0.20
SLIDE_W = 5143500
SLIDE_H = 9144000
SAFE_TOP_EMU = int(SLIDE_H * SAFE_TOP)
SAFE_BOTTOM_EMU = int(SLIDE_H * SAFE_BOTTOM)
CONTENT_TOP = SAFE_TOP_EMU
CONTENT_BOTTOM = SLIDE_H - SAFE_BOTTOM_EMU
CONTENT_HEIGHT = CONTENT_BOTTOM - CONTENT_TOP

# Add blank layout
blank_layout = prs.slide_layouts[6]  # Blank

# ========== LAYOUT 0: title_scene ==========
slide = prs.slides.add_slide(blank_layout)
# Title placeholder
left = Emu(int(SLIDE_W * 0.1))
top = Emu(int(CONTENT_TOP + CONTENT_HEIGHT * 0.15))
width = Emu(int(SLIDE_W * 0.8))
height = Emu(int(CONTENT_HEIGHT * 0.25))
title_box = slide.shapes.add_textbox(left, top, width, height)
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Title Headline"
p.font.size = Emu(480000)  # 48pt
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Subtitle
left = Emu(int(SLIDE_W * 0.1))
top = Emu(int(CONTENT_TOP + CONTENT_HEIGHT * 0.45))
width = Emu(int(SLIDE_W * 0.8))
height = Emu(int(CONTENT_HEIGHT * 0.15))
sub_box = slide.shapes.add_textbox(left, top, width, height)
tf = sub_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Subtitle text goes here"
p.font.size = Emu(240000)  # 24pt
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Background
background = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(SLIDE_W), Emu(SLIDE_H)
)
background.fill.solid()
background.fill.fore_color.rgb = DARK_BG
background.line.fill.background()
# Send to back
spTree = slide.shapes._spTree
sp = background._element
spTree.remove(sp)
spTree.insert(2, sp)

# ========== LAYOUT 1: stat_scene ==========
slide = prs.slides.add_slide(blank_layout)
# Big number
left = Emu(int(SLIDE_W * 0.1))
top = Emu(int(CONTENT_TOP + CONTENT_HEIGHT * 0.1))
width = Emu(int(SLIDE_W * 0.8))
height = Emu(int(CONTENT_HEIGHT * 0.35))
num_box = slide.shapes.add_textbox(left, top, width, height)
tf = num_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "$12M"
p.font.size = Emu(720000)  # 72pt
p.font.bold = True
p.font.color.rgb = PRIMARY_COLOR
p.alignment = PP_ALIGN.CENTER

# Label
left = Emu(int(SLIDE_W * 0.1))
top = Emu(int(CONTENT_TOP + CONTENT_HEIGHT * 0.50))
width = Emu(int(SLIDE_W * 0.8))
height = Emu(int(CONTENT_HEIGHT * 0.12))
label_box = slide.shapes.add_textbox(left, top, width, height)
tf = label_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Revenue"
p.font.size = Emu(360000)  # 36pt
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Sublabel
left = Emu(int(SLIDE_W * 0.1))
top = Emu(int(CONTENT_TOP + CONTENT_HEIGHT * 0.65))
width = Emu(int(SLIDE_W * 0.8))
height = Emu(int(CONTENT_HEIGHT * 0.10))
sub_box = slide.shapes.add_textbox(left, top, width, height)
tf = sub_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "+15% Year over Year"
p.font.size = Emu(240000)  # 24pt
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

background = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(SLIDE_W), Emu(SLIDE_H)
)
background.fill.solid()
background.fill.fore_color.rgb = DARK_BG
background.line.fill.background()
spTree = slide.shapes._spTree
sp = background._element
spTree.remove(sp)
spTree.insert(2, sp)

# ========== LAYOUT 2: bullet_scene ==========
slide = prs.slides.add_slide(blank_layout)
# Title
left = Emu(int(SLIDE_W * 0.1))
top = Emu(int(CONTENT_TOP + CONTENT_HEIGHT * 0.05))
width = Emu(int(SLIDE_W * 0.8))
height = Emu(int(CONTENT_HEIGHT * 0.15))
title_box = slide.shapes.add_textbox(left, top, width, height)
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Key Points"
p.font.size = Emu(400000)  # 40pt
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

# Bullet 1
left = Emu(int(SLIDE_W * 0.1))
top = Emu(int(CONTENT_TOP + CONTENT_HEIGHT * 0.25))
width = Emu(int(SLIDE_W * 0.8))
height = Emu(int(CONTENT_HEIGHT * 0.18))
b1 = slide.shapes.add_textbox(left, top, width, height)
tf = b1.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "• First bullet point here"
p.font.size = Emu(280000)  # 28pt
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

# Bullet 2
left = Emu(int(SLIDE_W * 0.1))
top = Emu(int(CONTENT_TOP + CONTENT_HEIGHT * 0.48))
width = Emu(int(SLIDE_W * 0.8))
height = Emu(int(CONTENT_HEIGHT * 0.18))
b2 = slide.shapes.add_textbox(left, top, width, height)
tf = b2.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "• Second bullet point here"
p.font.size = Emu(280000)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

# Bullet 3
left = Emu(int(SLIDE_W * 0.1))
top = Emu(int(CONTENT_TOP + CONTENT_HEIGHT * 0.71))
width = Emu(int(SLIDE_W * 0.8))
height = Emu(int(CONTENT_HEIGHT * 0.18))
b3 = slide.shapes.add_textbox(left, top, width, height)
tf = b3.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "• Third bullet point here"
p.font.size = Emu(280000)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

background = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(SLIDE_W), Emu(SLIDE_H)
)
background.fill.solid()
background.fill.fore_color.rgb = DARK_BG
background.line.fill.background()
spTree = slide.shapes._spTree
sp = background._element
spTree.remove(sp)
spTree.insert(2, sp)

# ========== LAYOUT 3: image_scene ==========
slide = prs.slides.add_slide(blank_layout)
# Image placeholder area
left = Emu(int(SLIDE_W * 0.05))
top = Emu(int(CONTENT_TOP + CONTENT_HEIGHT * 0.05))
width = Emu(int(SLIDE_W * 0.9))
height = Emu(int(CONTENT_HEIGHT * 0.7))
img_box = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, left, top, width, height
)
img_box.fill.solid()
img_box.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
img_box.line.color.rgb = PRIMARY_COLOR

# Caption
left = Emu(int(SLIDE_W * 0.1))
top = Emu(int(CONTENT_TOP + CONTENT_HEIGHT * 0.80))
width = Emu(int(SLIDE_W * 0.8))
height = Emu(int(CONTENT_HEIGHT * 0.12))
cap_box = slide.shapes.add_textbox(left, top, width, height)
tf = cap_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Image Caption"
p.font.size = Emu(280000)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

background = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(SLIDE_W), Emu(SLIDE_H)
)
background.fill.solid()
background.fill.fore_color.rgb = DARK_BG
background.line.fill.background()
spTree = slide.shapes._spTree
sp = background._element
spTree.remove(sp)
spTree.insert(2, sp)

# ========== LAYOUT 4: cta_scene ==========
slide = prs.slides.add_slide(blank_layout)
# CTA Headline
left = Emu(int(SLIDE_W * 0.1))
top = Emu(int(CONTENT_TOP + CONTENT_HEIGHT * 0.2))
width = Emu(int(SLIDE_W * 0.8))
height = Emu(int(CONTENT_HEIGHT * 0.25))
cta_box = slide.shapes.add_textbox(left, top, width, height)
tf = cta_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Call to Action"
p.font.size = Emu(480000)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Subheadline
left = Emu(int(SLIDE_W * 0.1))
top = Emu(int(CONTENT_TOP + CONTENT_HEIGHT * 0.50))
width = Emu(int(SLIDE_W * 0.8))
height = Emu(int(CONTENT_HEIGHT * 0.15))
sub_box = slide.shapes.add_textbox(left, top, width, height)
tf = sub_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Follow for more insights"
p.font.size = Emu(280000)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Arrow button placeholder
left = Emu(int(SLIDE_W * 0.3))
top = Emu(int(CONTENT_TOP + CONTENT_HEIGHT * 0.70))
width = Emu(int(SLIDE_W * 0.4))
height = Emu(int(CONTENT_HEIGHT * 0.12))
btn = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
)
btn.fill.solid()
btn.fill.fore_color.rgb = PRIMARY_COLOR
btn.line.fill.background()

background = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(SLIDE_W), Emu(SLIDE_H)
)
background.fill.solid()
background.fill.fore_color.rgb = DARK_BG
background.line.fill.background()
spTree = slide.shapes._spTree
sp = background._element
spTree.remove(sp)
spTree.insert(2, sp)

# Save
output_path = "python/reel_converter/templates/reel_clean.pptx"
prs.save(output_path)
print(f"Created: {output_path}")
print(f"Dimensions: {prs.slide_width} x {prs.slide_height} EMU (5.625\" x 10\")")
print(f"Slides: {len(prs.slides)}")
