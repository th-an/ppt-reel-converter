from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# Instagram Reels 9:16 aspect ratio
# 5.625" x 10" (13.333" x 7.5" is landscape 16:9)
# For 9:16 portrait, we need width = 5.625, height = 10
# 1 inch = 914400 EMU
SLIDE_W = Emu(5143500)  # 5.625"
SLIDE_H = Emu(9144000)  # 10"

# Safe zones
SAFE_TOP = 0.15
SAFE_BOTTOM = 0.20
CONTENT_TOP = int(SLIDE_H * SAFE_TOP)
CONTENT_BOTTOM = int(SLIDE_H * (1 - SAFE_BOTTOM))
CONTENT_HEIGHT = CONTENT_BOTTOM - CONTENT_TOP

# Colors
PRIMARY = RGBColor(0x01, 0x96, 0xFF)  # #0196FF
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Clear default slide
blank_layout = prs.slide_layouts[6]  # Blank

# Get the slide master
slide_master = prs.slide_masters[0]

# ========== LAYOUT 0: title_scene (Title Only) ==========
# Use the "Title Only" layout (index 5)
title_only_layout = prs.slide_layouts[5]

# Add a slide to show the layout
slide = prs.slides.add_slide(title_only_layout)
# Set background
background = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H
)
background.fill.solid()
background.fill.fore_color.rgb = DARK_BG
background.line.fill.background()

# Add title text
for shape in slide.shapes:
    if shape.is_placeholder:
        ph = shape.placeholder_format
        if ph.type == 1:  # TITLE
            shape.text_frame.text = "Title Headline"
            p = shape.text_frame.paragraphs[0]
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

# ========== LAYOUT 1: stat_scene (Title + Content) ==========
title_content_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(title_content_layout)

# Background
background = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H
)
background.fill.solid()
background.fill.fore_color.rgb = DARK_BG
background.line.fill.background()

# Fill placeholders
for shape in slide.shapes:
    if shape.is_placeholder:
        ph = shape.placeholder_format
        if ph.type == 1:  # TITLE
            shape.text_frame.text = "$12M"
            p = shape.text_frame.paragraphs[0]
            p.font.size = Pt(72)
            p.font.bold = True
            p.font.color.rgb = PRIMARY
            p.alignment = PP_ALIGN.CENTER
        elif ph.type == 7:  # OBJECT (content)
            shape.text_frame.text = "Revenue\n+15% Year over Year"
            p = shape.text_frame.paragraphs[0]
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            # Second paragraph
            if len(shape.text_frame.paragraphs) > 1:
                shape.text_frame.paragraphs[1].font.size = Pt(24)
                shape.text_frame.paragraphs[1].font.color.rgb = WHITE

# ========== LAYOUT 2: bullet_scene (Title + Content) ==========
slide = prs.slides.add_slide(title_content_layout)

# Background
background = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H
)
background.fill.solid()
background.fill.fore_color.rgb = DARK_BG
background.line.fill.background()

for shape in slide.shapes:
    if shape.is_placeholder:
        ph = shape.placeholder_format
        if ph.type == 1:  # TITLE
            shape.text_frame.text = "Key Points"
            p = shape.text_frame.paragraphs[0]
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.LEFT
        elif ph.type == 7:  # OBJECT (content)
            tf = shape.text_frame
            tf.text = "First bullet point"
            p = tf.paragraphs[0]
            p.font.size = Pt(28)
            p.font.color.rgb = WHITE
            p.level = 0
            
            p = tf.add_paragraph()
            p.text = "Second bullet point"
            p.font.size = Pt(28)
            p.font.color.rgb = WHITE
            p.level = 0
            
            p = tf.add_paragraph()
            p.text = "Third bullet point"
            p.font.size = Pt(28)
            p.font.color.rgb = WHITE
            p.level = 0

# ========== LAYOUT 3: image_scene (Picture with Caption) ==========
picture_caption_layout = prs.slide_layouts[8]
slide = prs.slides.add_slide(picture_caption_layout)

# Background
background = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H
)
background.fill.solid()
background.fill.fore_color.rgb = DARK_BG
background.line.fill.background()

for shape in slide.shapes:
    if shape.is_placeholder:
        ph = shape.placeholder_format
        if ph.type == 1:  # TITLE
            shape.text_frame.text = "Image Caption"
            p = shape.text_frame.paragraphs[0]
            p.font.size = Pt(28)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        elif ph.type == 18:  # PICTURE
            # Picture placeholder - add a placeholder shape
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)

# ========== LAYOUT 4: cta_scene (Title Only) ==========
slide = prs.slides.add_slide(title_only_layout)

# Background
background = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H
)
background.fill.solid()
background.fill.fore_color.rgb = PRIMARY
background.line.fill.background()

for shape in slide.shapes:
    if shape.is_placeholder:
        ph = shape.placeholder_format
        if ph.type == 1:  # TITLE
            shape.text_frame.text = "Call to Action"
            p = shape.text_frame.paragraphs[0]
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

# Save
output_path = "python/reel_converter/templates/reel_clean.pptx"
prs.save(output_path)
print(f"Created: {output_path}")
print(f"Dimensions: {prs.slide_width} x {prs.slide_height} EMU")
print(f"Slides: {len(prs.slides)}")
print(f"Layouts: {len(prs.slide_layouts)}")

# Now profile it
import sys
sys.path.insert(0, "python")
from reel_converter.profile import profile_template, generate_config

config = profile_template(output_path)
print(f"\nProfiled template: {config.template_name}")
print(f"Dimensions: {config.dimensions.width_emu} x {config.dimensions.height_emu}")
print(f"Layouts: {list(config.layout_mappings.keys())}")
print(f"Routing: {config.content_type_routing}")

# Save config
config_path, catalog_path = generate_config(config, "python/reel_converter/templates")
print(f"\nConfig saved: {config_path}")
print(f"Catalog saved: {catalog_path}")
