import sys
sys.path.insert(0, "python")

from pptx import Presentation

for slide_num in [1, 2]:
    path = f"output/slide_{slide_num}_scenes.pptx"
    prs = Presentation(path)
    
    print(f"\n=== {path} ===")
    print(f"Slides: {len(prs.slides)}")
    
    for idx, slide in enumerate(prs.slides):
        print(f"\n  Slide {idx + 1}: Layout={slide.slide_layout.name}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    print(f"    Shape '{shape.name}' (placeholder={shape.is_placeholder}):")
                    print(f"      Text: '{text[:80]}'")
                    # Check font info (paragraph-level and run-level)
                    for para in shape.text_frame.paragraphs:
                        para_font_size = para.font.size.pt if para.font.size else "default"
                        try:
                            para_color = para.font.color.rgb
                        except AttributeError:
                            para_color = "default"
                        print(f"      Para: {para_font_size}pt, bold={para.font.bold}, color={para_color}")
                        for run in para.runs:
                            font_size = run.font.size.pt if run.font.size else "default"
                            try:
                                color = run.font.color.rgb
                            except AttributeError:
                                color = "default"
                            print(f"      Run: {font_size}pt, bold={run.font.bold}, color={color}")
                else:
                    print(f"    Shape '{shape.name}': (empty)")
            else:
                print(f"    Shape '{shape.name}': (no text frame)")
