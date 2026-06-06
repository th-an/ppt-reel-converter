import sys
sys.path.insert(0, "python")

from pptx import Presentation

path = "output_complex/combined_reel.pptx"
prs = Presentation(path)

print(f"=== Combined Reel (Complex Test) ===")
print(f"Slides: {len(prs.slides)}")
print(f"Dimensions: {prs.slide_width} x {prs.slide_height}")

for idx, slide in enumerate(prs.slides):
    print(f"\n  Slide {idx + 1}: {slide.slide_layout.name}")
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            text = shape.text_frame.text.strip()[:80]
            print(f"    - {shape.name}: '{text}'")

# Also check summary
import json
with open("output_complex/summary.json") as f:
    summary = json.load(f)

print(f"\n=== Summary ===")
print(f"Total slides: {summary['total_slides']}")
print(f"Total scenes: {summary['total_scenes']}")
for result in summary['slide_results']:
    print(f"  Slide {result['slide']}: {result['scenes']} scenes, passed={result['passed']}")
