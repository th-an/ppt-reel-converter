import sys
sys.path.insert(0, "python")
from pptx import Presentation

prs = Presentation("python/reel_converter/templates/reel_clean.pptx")
print(f"Template: {len(prs.slides)} slides")
print(f"Layouts: {len(prs.slide_layouts)}")

for idx, layout in enumerate(prs.slide_layouts):
    print(f"\nLayout {idx}: {layout.name}")
    for shape in layout.placeholders:
        print(f"  Placeholder idx={shape.placeholder_format.idx}, type={shape.placeholder_format.type}, name={shape.name}")

print("\n--- Slide 1 (Title Slide) ---")
for shape in prs.slides[0].shapes:
    print(f"  Shape: {shape.name}, is_placeholder={shape.is_placeholder}, placeholder_type={shape.placeholder_format.type if shape.is_placeholder else 'N/A'}")
    if shape.has_text_frame:
        print(f"    Text: '{shape.text_frame.text.strip()}'")
