import sys
sys.path.insert(0, "python")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import io

# Create a simple test image
img = Image.new('RGB', (400, 300), color='steelblue')
img_buffer = io.BytesIO()
img.save(img_buffer, format='PNG')
img_buffer.seek(0)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title with image
blank_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(blank_layout)

# Title
left = Inches(0.5)
top = Inches(0.5)
width = Inches(9)
height = Inches(1)
title_box = slide1.shapes.add_textbox(left, top, width, height)
tf = title_box.text_frame
tf.text = "Product Features"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

# Add image
left = Inches(2)
top = Inches(2)
width = Inches(6)
height = Inches(4)
slide1.shapes.add_picture(img_buffer, left, top, width, height)

# Slide 2: Title with bullets
slide2 = prs.slides.add_slide(blank_layout)

# Title
left = Inches(0.5)
top = Inches(0.5)
width = Inches(9)
height = Inches(1)
title_box = slide2.shapes.add_textbox(left, top, width, height)
tf = title_box.text_frame
tf.text = "Key Metrics"
p = tf.paragraphs[0]
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

# Bullets
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(4)
body_box = slide2.shapes.add_textbox(left, top, width, height)
tf = body_box.text_frame
tf.word_wrap = True

items = [
    "Revenue: $12M (up 40%)",
    "Users: 72,000 (active)",
    "Conversion: 4.8% (avg)",
    "Churn: 2.1% (low)",
]

for i, item in enumerate(items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = f"• {item}"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.space_after = Pt(12)

# Save
output_path = "tests/fixtures/image_test.pptx"
prs.save(output_path)
print(f"Created {output_path} with {len(prs.slides)} slides")
print("Slide 1: Product Features + image")
print("Slide 2: Key Metrics + bullets")
