import sys
sys.path.insert(0, "python")

from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Template configurations with different color schemes and styles
TEMPLATES = {
    "modern": {
        "bg_color": "F5F5F7",
        "accent_color": "FF6B35",
        "text_color": "1A1A2E",
        "subtitle_color": "6B7280",
        "font_header": "Inter",
        "font_body": "Inter",
        "style": "modern"
    },
    "bold": {
        "bg_color": "FF3366",
        "accent_color": "FFFFFF",
        "text_color": "FFFFFF",
        "subtitle_color": "FFE4E1",
        "font_header": "Montserrat",
        "font_body": "Montserrat",
        "style": "bold"
    },
    "minimal": {
        "bg_color": "FFFFFF",
        "accent_color": "000000",
        "text_color": "000000",
        "subtitle_color": "666666",
        "font_header": "Helvetica Neue",
        "font_body": "Helvetica Neue",
        "style": "minimal"
    },
    "corporate": {
        "bg_color": "1E3A5F",
        "accent_color": "D4AF37",
        "text_color": "FFFFFF",
        "subtitle_color": "B0C4DE",
        "font_header": "Georgia",
        "font_body": "Arial",
        "style": "corporate"
    }
}

def create_template(name, config):
    """Create a 9:16 PPTX template with the given color scheme."""
    prs = Presentation()
    prs.slide_width = Emu(5143500)   # 5.625 inches
    prs.slide_height = Emu(9144000)  # 10 inches
    
    # Parse colors
    bg_color = RGBColor.from_string(config["bg_color"])
    accent_color = RGBColor.from_string(config["accent_color"])
    text_color = RGBColor.from_string(config["text_color"])
    subtitle_color = RGBColor.from_string(config["subtitle_color"])
    
    # Layout 0: Title Slide
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Add background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    
    # Add title placeholder
    title = slide.shapes.add_textbox(Emu(200000), Emu(2000000), Emu(4743500), Emu(1000000))
    tf = title.text_frame
    tf.text = "Title"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = config["font_header"]
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle placeholder
    subtitle = slide.shapes.add_textbox(Emu(200000), Emu(3200000), Emu(4743500), Emu(600000))
    tf = subtitle.text_frame
    tf.text = "Subtitle"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = subtitle_color
    p.font.name = config["font_body"]
    p.alignment = PP_ALIGN.CENTER
    
    # Layout 1: Title and Content
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    
    title = slide.shapes.add_textbox(Emu(200000), Emu(500000), Emu(4743500), Emu(800000))
    tf = title.text_frame
    tf.text = "Title"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = config["font_header"]
    p.alignment = PP_ALIGN.LEFT
    
    body = slide.shapes.add_textbox(Emu(200000), Emu(1500000), Emu(4743500), Emu(5000000))
    tf = body.text_frame
    tf.word_wrap = True
    tf.text = "Body text"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = text_color
    p.font.name = config["font_body"]
    
    # Layout 2: Section Header
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    
    title = slide.shapes.add_textbox(Emu(200000), Emu(3000000), Emu(4743500), Emu(1000000))
    tf = title.text_frame
    tf.text = "Section"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p.font.name = config["font_header"]
    p.alignment = PP_ALIGN.CENTER
    
    # Layout 3: Picture with Caption
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    
    # Image placeholder
    img = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(200000), Emu(1000000), Emu(4743500), Emu(4000000))
    img.fill.solid()
    img.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    img.line.color.rgb = RGBColor(0x99, 0x99, 0x99)
    
    caption = slide.shapes.add_textbox(Emu(200000), Emu(5200000), Emu(4743500), Emu(600000))
    tf = caption.text_frame
    tf.text = "Caption"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = text_color
    p.font.name = config["font_body"]
    p.alignment = PP_ALIGN.CENTER
    
    # Layout 4: CTA
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = accent_color
    bg.line.fill.background()
    
    title = slide.shapes.add_textbox(Emu(200000), Emu(2500000), Emu(4743500), Emu(1000000))
    tf = title.text_frame
    tf.text = "Call to Action"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.name = config["font_header"]
    p.alignment = PP_ALIGN.CENTER
    
    # Save
    path = f"python/reel_converter/templates/reel_{name}.pptx"
    prs.save(path)
    print(f"Created {path}")
    return path

# Create all templates
for name, config in TEMPLATES.items():
    create_template(name, config)

print("\nAll templates created!")
