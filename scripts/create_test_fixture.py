from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide 1: Title only
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Q3 Revenue Report"

# Slide 2: Bullets with image
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Key Metrics"
body = slide.shapes.placeholders[1]
tf = body.text_frame
tf.text = "Revenue: $12M"
p = tf.add_paragraph()
p.text = "Operating margin: 34%"
p.level = 0
p = tf.add_paragraph()
p.text = "New customers: 2,400"
p.level = 0
p = tf.add_paragraph()
p.text = "Churn: 2.1%"
p.level = 0

prs.save("tests/fixtures/simple_test.pptx")
print("Created test fixture: tests/fixtures/simple_test.pptx")
