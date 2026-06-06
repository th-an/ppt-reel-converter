import sys
sys.path.insert(0, "python")

from pptx import Presentation

path = "output/combined_reel.pptx"
prs = Presentation(path)

print(f"=== Combined Reel ===")
print(f"Slides: {len(prs.slides)}")
print(f"Dimensions: {prs.slide_width} x {prs.slide_height}")

for idx, slide in enumerate(prs.slides):
    print(f"\n  Slide {idx + 1}: {slide.slide_layout.name}")
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            text = shape.text_frame.text.strip()[:60]
            print(f"    - {shape.name}: '{text}'")
