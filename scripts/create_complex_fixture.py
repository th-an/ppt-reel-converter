import sys
sys.path.insert(0, "python")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title with subtitle
blank_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(blank_layout)

# Title
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9)
height = Inches(1)
title_box = slide1.shapes.add_textbox(left, top, width, height)
tf = title_box.text_frame
tf.text = "Q3 Product Launch Report"
p = tf.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
p.alignment = PP_ALIGN.CENTER

# Subtitle
left = Inches(0.5)
top = Inches(2.5)
width = Inches(9)
height = Inches(0.5)
sub_box = slide1.shapes.add_textbox(left, top, width, height)
tf = sub_box.text_frame
tf.text = "Prepared by Marketing Team — October 2025"
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
p.alignment = PP_ALIGN.CENTER

# Slide 2: Bullets with image
slide2 = prs.slides.add_slide(blank_layout)

# Title
left = Inches(0.5)
top = Inches(0.5)
width = Inches(9)
height = Inches(0.8)
title_box = slide2.shapes.add_textbox(left, top, width, height)
tf = title_box.text_frame
tf.text = "Key Features"
p = tf.paragraphs[0]
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

# Bullets (left side)
left = Inches(0.5)
top = Inches(1.5)
width = Inches(4.5)
height = Inches(4)
body_box = slide2.shapes.add_textbox(left, top, width, height)
tf = body_box.text_frame
tf.word_wrap = True

items = [
    "AI-powered recommendations",
    "Real-time collaboration",
    "Dark mode support",
    "Cross-platform sync",
    "Enterprise SSO",
    "Custom branding"
]

for i, item in enumerate(items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = f"• {item}"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.space_after = Pt(12)

# Image placeholder (right side)
left = Inches(5.5)
top = Inches(1.5)
width = Inches(4)
height = Inches(4)
img_shape = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
img_shape.fill.solid()
img_shape.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
img_shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

# Add "Image Placeholder" text
img_tf = img_shape.text_frame
img_tf.text = "[Product Screenshot]"
p = img_tf.paragraphs[0]
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
p.alignment = PP_ALIGN.CENTER

# Slide 3: Table with stats
slide3 = prs.slides.add_slide(blank_layout)

# Title
left = Inches(0.5)
top = Inches(0.5)
width = Inches(9)
height = Inches(0.8)
title_box = slide3.shapes.add_textbox(left, top, width, height)
tf = title_box.text_frame
tf.text = "Performance Metrics"
p = tf.paragraphs[0]
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

# Table
rows = 4
cols = 3
left = Inches(1.5)
top = Inches(1.5)
width = Inches(7)
height = Inches(3)
table = slide3.shapes.add_table(rows, cols, left, top, width, height).table

# Set column headers
headers = ["Metric", "Q2", "Q3"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(18)
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x01, 0x96, 0xFF)

# Data
data = [
    ["Revenue", "$8.5M", "$12.3M"],
    ["Users", "45,000", "72,000"],
    ["Conversion", "3.2%", "4.8%"],
]

for row_idx, row_data in enumerate(data, start=1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Slide 4: Chart (represented as text for now)
slide4 = prs.slides.add_slide(blank_layout)

# Title
left = Inches(0.5)
top = Inches(0.5)
width = Inches(9)
height = Inches(0.8)
title_box = slide4.shapes.add_textbox(left, top, width, height)
tf = title_box.text_frame
tf.text = "Growth Trend"
p = tf.paragraphs[0]
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

# Chart placeholder
left = Inches(1.5)
top = Inches(1.5)
width = Inches(7)
height = Inches(4)
chart_shape = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
chart_shape.fill.solid()
chart_shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
chart_shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

# Add chart text
chart_tf = chart_shape.text_frame
chart_tf.text = "[Chart: Revenue Growth 2023-2025]\n\nJan: $2M → Jun: $5M → Oct: $12M"
p = chart_tf.paragraphs[0]
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
p.alignment = PP_ALIGN.CENTER

# Slide 5: CTA / Summary
slide5 = prs.slides.add_slide(blank_layout)

# Title
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9)
height = Inches(1)
title_box = slide5.shapes.add_textbox(left, top, width, height)
tf = title_box.text_frame
tf.text = "Ready to Scale"
p = tf.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
p.alignment = PP_ALIGN.CENTER

# Subtitle
left = Inches(0.5)
top = Inches(2.5)
width = Inches(9)
height = Inches(0.5)
sub_box = slide5.shapes.add_textbox(left, top, width, height)
tf = sub_box.text_frame
tf.text = "Contact sales@company.com to get started"
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
p.alignment = PP_ALIGN.CENTER

# Save
output_path = "tests/fixtures/complex_test.pptx"
prs.save(output_path)
print(f"Created {output_path} with {len(prs.slides)} slides")
print("Slides:")
for i, slide in enumerate(prs.slides, 1):
    shapes = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
    print(f"  Slide {i}: {len(shapes)} text shapes")
